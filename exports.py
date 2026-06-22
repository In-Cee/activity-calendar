# ============================================================
#  Export helpers - v0.6
#  Build PowerPoint and Excel exports for the Exec Brief
# ============================================================

import io
from datetime import date
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


FOUNDATION_ORANGE_RGB = RGBColor(0xF3, 0x70, 0x21)
DARK_GREY_RGB = RGBColor(0x2B, 0x2B, 0x2B)
LIGHT_GREY_RGB = RGBColor(0x88, 0x88, 0x88)


def build_pptx(narrative, kpis, attention_df, top_func, top_country):
    """Build a single-slide Foundation-branded PowerPoint Exec Brief."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Orange top band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.6))
    band.fill.solid()
    band.fill.fore_color.rgb = FOUNDATION_ORANGE_RGB
    band.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.5))
    tf = title.text_frame
    tf.text = "Activity Calendar - Executive Brief"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Date sub-line
    sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.3))
    sub.text_frame.text = f"Mastercard Foundation · Enterprise Planning · {date.today().strftime('%d %B %Y')}"
    sub.text_frame.paragraphs[0].font.size = Pt(11)
    sub.text_frame.paragraphs[0].font.color.rgb = LIGHT_GREY_RGB

    # Hero narrative
    hero = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(12.5), Inches(1.4))
    hf = hero.text_frame
    hf.word_wrap = True
    p = hf.paragraphs[0]
    p.text = "What you need to know"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = FOUNDATION_ORANGE_RGB

    p2 = hf.add_paragraph()
    p2.text = narrative
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GREY_RGB

    # KPI strip (5 boxes)
    box_w = Inches(2.4)
    box_h = Inches(0.9)
    box_top = Inches(2.7)
    box_left_start = Inches(0.4)
    gap = Inches(0.15)
    for i, (label, val) in enumerate(kpis):
        left = box_left_start + (box_w + gap) * i
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, box_top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        box.line.color.rgb = FOUNDATION_ORANGE_RGB
        box.line.width = Pt(1.5)

        txt = slide.shapes.add_textbox(left, box_top + Inches(0.1), box_w, Inches(0.35))
        txt.text_frame.text = label
        txt.text_frame.paragraphs[0].font.size = Pt(10)
        txt.text_frame.paragraphs[0].font.color.rgb = LIGHT_GREY_RGB

        val_box = slide.shapes.add_textbox(left, box_top + Inches(0.4), box_w, Inches(0.5))
        val_box.text_frame.text = str(val)
        val_box.text_frame.paragraphs[0].font.size = Pt(22)
        val_box.text_frame.paragraphs[0].font.bold = True
        val_box.text_frame.paragraphs[0].font.color.rgb = DARK_GREY_RGB

    # Functions needing attention
    att_box = slide.shapes.add_textbox(Inches(0.4), Inches(3.9), Inches(12.5), Inches(0.4))
    att_box.text_frame.text = "Functions needing attention"
    att_box.text_frame.paragraphs[0].font.size = Pt(14)
    att_box.text_frame.paragraphs[0].font.bold = True
    att_box.text_frame.paragraphs[0].font.color.rgb = FOUNDATION_ORANGE_RGB

    if attention_df is None or attention_df.empty:
        ok_box = slide.shapes.add_textbox(Inches(0.4), Inches(4.4), Inches(12.5), Inches(0.4))
        ok_box.text_frame.text = "✓ All functions are within comfort thresholds. No immediate action needed."
        ok_box.text_frame.paragraphs[0].font.size = Pt(12)
        ok_box.text_frame.paragraphs[0].font.color.rgb = DARK_GREY_RGB
    else:
        rows = min(5, len(attention_df))
        cols = 5
        table_shape = slide.shapes.add_table(rows + 1, cols,
                                             Inches(0.4), Inches(4.4),
                                             Inches(12.5), Inches(0.5 + rows * 0.35))
        table = table_shape.table
        for i, h in enumerate(["Function", "Peak week", "Status", "Elevated", "Critical"]):
            cell = table.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = FOUNDATION_ORANGE_RGB
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for r in range(rows):
            row = attention_df.iloc[r]
            for c, val in enumerate([row["Function"], row["Peak week"], row["Status"],
                                     row["Elevated"], row["Critical"]]):
                cell = table.cell(r + 1, c)
                cell.text = str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = DARK_GREY_RGB

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3))
    footer.text_frame.text = (
        f"Top hosting function: {top_func} · Top country: {top_country} · "
        f"Auto-generated by Activity Calendar"
    )
    footer.text_frame.paragraphs[0].font.size = Pt(9)
    footer.text_frame.paragraphs[0].font.color.rgb = LIGHT_GREY_RGB
    footer.text_frame.paragraphs[0].font.italic = True

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_excel_export(view_df):
    """Return a BytesIO buffer of view_df formatted as a Foundation Excel file."""
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        view_df.to_excel(writer, sheet_name="Activities", index=False, startrow=1)
        workbook = writer.book
        sheet = writer.sheets["Activities"]

        header_fmt = workbook.add_format({
            "bold": True, "bg_color": "#F37021", "font_color": "white",
            "align": "center", "valign": "vcenter", "border": 1,
        })
        title_fmt = workbook.add_format({
            "bold": True, "font_size": 14, "font_color": "#F37021",
        })
        sheet.write(0, 0, f"Activity Calendar export · {date.today().strftime('%d %B %Y')}", title_fmt)
        for col_num, value in enumerate(view_df.columns.values):
            sheet.write(1, col_num, value, header_fmt)
        for i, col in enumerate(view_df.columns):
            try:
                max_len = max(view_df[col].astype(str).map(len).max(), len(str(col))) + 2
            except Exception:
                max_len = 20
            sheet.set_column(i, i, min(max_len, 50))
    buf.seek(0)
    return buf


def build_template_excel():
    """Return an empty mass-upload template with headers."""
    headers = [
        "Activity ID (Auto)", "StartDate", "EndDate", "Type", "Title", "Location",
        "Initiating Function", "Initiating Sub-Function", "Attendee Category",
        "Participating Function", "Participating Sub-Function", "Internal/External",
        "Month (Auto/Manual)", "Year (Auto/Manual)", "Note"
    ]
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        df = pd.DataFrame(columns=headers)
        df.to_excel(writer, sheet_name="Activities", index=False, startrow=1)
        workbook = writer.book
        sheet = writer.sheets["Activities"]
        guide_fmt = workbook.add_format({
            "bold": True, "bg_color": "#FFF4CC", "align": "left",
        })
        header_fmt = workbook.add_format({
            "bold": True, "bg_color": "#E8E8E8", "align": "center", "border": 1,
        })
        sheet.merge_range(0, 0, 0, len(headers) - 1,
                          "INPUT GUIDE: Fill columns B–O for each activity. Activity ID is auto-generated on upload.",
                          guide_fmt)
        for i, h in enumerate(headers):
            sheet.write(1, i, h, header_fmt)
            sheet.set_column(i, i, max(16, len(h) + 2))
    buf.seek(0)
    return buf
